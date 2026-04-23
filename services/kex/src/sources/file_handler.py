"""
File Text Extractor for KEX Service
Extracts plain text from common document formats.
Supported: PDF, DOCX, CSV, JSON, XML, plain text.
"""

import csv
import io
import json
import logging

logger = logging.getLogger(__name__)


def extract_text(file_bytes: bytes, mimetype: str, filename: str = "document") -> str:
    """
    Extract textual content from file bytes.

    Args:
        file_bytes: Raw file content.
        mimetype:   MIME type string (e.g. "application/pdf").
        filename:   Original filename (used for format detection).

    Returns:
        Extracted text as a single string.

    Raises:
        ValueError: If the mimetype is unsupported or extraction fails.
    """
    # Try unstructured.io first for enhanced formats (OCR, images, etc.)
    try:
        from .unstructured_handler import should_use_unstructured, extract_text as unstructured_extract
        if should_use_unstructured(mimetype):
            logger.info(f"Using unstructured.io for {mimetype}")
            return unstructured_extract(file_bytes, mimetype, filename)
    except Exception as exc:
        logger.warning(f"Unstructured.io fallback: {exc}, using built-in extractor")

    # Normalise mimetype (strip parameters like charset)
    base_mime = mimetype.split(";")[0].strip().lower()

    if base_mime == "application/pdf":
        return _extract_pdf(file_bytes)

    if base_mime in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/docx",
    ):
        return _extract_docx(file_bytes)

    if base_mime in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return _extract_pptx(file_bytes)

    if base_mime in ("text/csv", "application/csv"):
        return _extract_csv(file_bytes)

    if base_mime in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return _extract_xlsx(file_bytes)

    if base_mime in ("application/json",):
        return _extract_json(file_bytes)

    if base_mime in ("text/xml", "application/xml"):
        return _extract_xml(file_bytes)

    if base_mime == "application/rtf":
        return _extract_rtf(file_bytes)

    if base_mime == "application/epub+zip":
        return _extract_epub(file_bytes)

    if base_mime in ("message/rfc822", "application/eml"):
        return _extract_eml(file_bytes)

    if base_mime == "application/vnd.ms-outlook":
        return _extract_msg(file_bytes)

    if base_mime == "application/vnd.oasis.opendocument.text":
        return _extract_odt(file_bytes)

    if base_mime in ("application/x-yaml", "text/yaml"):
        return _extract_yaml(file_bytes)

    if base_mime == "application/toml":
        return _extract_toml(file_bytes)

    if base_mime in ("image/png", "image/jpeg", "image/tiff", "image/bmp", "image/gif", "image/webp"):
        return _extract_image_ocr(file_bytes)

    if base_mime in ("audio/mpeg", "audio/wav", "audio/ogg", "audio/flac", "audio/mp4", "audio/webm"):
        return _extract_audio(file_bytes, base_mime)

    if base_mime in ("video/mp4", "video/webm", "video/quicktime", "video/x-msvideo", "video/x-matroska", "video/mpeg"):
        return _extract_video(file_bytes, base_mime)

    if base_mime.startswith("text/"):
        # text/plain, text/html, text/markdown, etc.
        return _extract_plaintext(file_bytes)

    raise ValueError(f"Unsupported mimetype: {mimetype}")


# ── format-specific extractors ────────────────────────────────────────


def _extract_pdf(data: bytes) -> str:
    """Extract text from all pages of a PDF."""
    try:
        import PyPDF2  # type: ignore
    except ImportError:
        raise ValueError("PyPDF2 is not installed")

    reader = PyPDF2.PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page_num, page in enumerate(reader.pages):
        try:
            text = page.extract_text()
            if text:
                parts.append(text.strip())
        except Exception as exc:
            logger.warning(f"PDF page {page_num} extraction failed: {exc}")

    if not parts:
        raise ValueError("PDF contained no extractable text")

    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    """Extract text from all slides of a PPTX file."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise ValueError("python-pptx is not installed")

    prs = Presentation(io.BytesIO(data))
    texts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        slide_texts.append(" | ".join(row_texts))
        if slide_texts:
            texts.append(f"Slide {slide_num}:\n" + "\n".join(slide_texts))

    if not texts:
        raise ValueError("PPTX contained no extractable text")

    return "\n\n".join(texts)


def _extract_docx(data: bytes) -> str:
    """Extract text from all paragraphs of a DOCX file."""
    try:
        from docx import Document  # type: ignore
    except ImportError:
        raise ValueError("python-docx is not installed")

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    if not paragraphs:
        raise ValueError("DOCX contained no extractable text")

    return "\n\n".join(paragraphs)


def _extract_csv(data: bytes) -> str:
    """Convert CSV to a human-readable text representation."""
    # Try UTF-8, fall back to latin-1
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            text = data.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        raise ValueError("Could not decode CSV with supported encodings")

    reader = csv.DictReader(io.StringIO(text))
    lines: list[str] = []

    fieldnames = reader.fieldnames or []
    if fieldnames:
        lines.append("Columns: " + ", ".join(str(f) for f in fieldnames))
        lines.append("")

    for i, row in enumerate(reader):
        row_parts = [f"{k}: {v}" for k, v in row.items() if v]
        if row_parts:
            lines.append(f"Row {i + 1}: " + " | ".join(row_parts))

    if not lines:
        raise ValueError("CSV contained no data")

    return "\n".join(lines)


def _extract_json(data: bytes) -> str:
    """
    Extract readable text from JSON.
    Tries to pretty-print structured data and also harvests string values.
    """
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            text = data.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        raise ValueError("Could not decode JSON")

    try:
        obj = json.loads(text)
    except json.JSONDecodeError as exc:
        raise ValueError(f"Invalid JSON: {exc}")

    # Harvest string leaves for NER to have clean prose
    strings: list[str] = []
    _collect_strings(obj, strings)

    if strings:
        return "\n".join(strings)

    # Fallback: pretty-printed JSON
    return json.dumps(obj, indent=2, ensure_ascii=False)


def _collect_strings(obj, out: list[str]) -> None:
    """Recursively collect string values from a JSON object."""
    if isinstance(obj, str):
        if len(obj.strip()) > 3:
            out.append(obj.strip())
    elif isinstance(obj, dict):
        for v in obj.values():
            _collect_strings(v, out)
    elif isinstance(obj, list):
        for item in obj:
            _collect_strings(item, out)


def _extract_xml(data: bytes) -> str:
    """Extract text content from XML by stripping tags."""
    try:
        from bs4 import BeautifulSoup  # type: ignore
        soup = BeautifulSoup(data, "xml")
        text = soup.get_text(separator="\n")
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if not lines:
            raise ValueError("XML contained no text content")
        return "\n".join(lines)
    except ImportError:
        # Fallback: basic tag stripping without bs4
        import re
        text = data.decode("utf-8", errors="replace")
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
        if not text:
            raise ValueError("XML contained no text content")
        return text


def _extract_plaintext(data: bytes) -> str:
    """Decode plain text with best-effort encoding detection."""
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


# ── XLSX (Excel) ───────────────────────────────────────────────────


def _extract_xlsx(data: bytes) -> str:
    """Extract text from all sheets of an Excel file."""
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        raise ValueError("openpyxl is not installed")

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"Sheet: {sheet_name}\n" + "\n".join(rows))

    wb.close()
    if not parts:
        raise ValueError("Excel file contained no data")
    return "\n\n".join(parts)


# ── RTF ────────────────────────────────────────────────────────────


def _extract_rtf(data: bytes) -> str:
    """Extract plain text from an RTF file."""
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
    except ImportError:
        raise ValueError("striprtf is not installed")

    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            rtf_content = data.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        rtf_content = data.decode("utf-8", errors="replace")

    text = rtf_to_text(rtf_content)
    if not text.strip():
        raise ValueError("RTF contained no extractable text")
    return text.strip()


# ── EPUB ───────────────────────────────────────────────────────────


def _extract_epub(data: bytes) -> str:
    """Extract text from an EPUB e-book."""
    try:
        import ebooklib  # type: ignore
        from ebooklib import epub  # type: ignore
        from bs4 import BeautifulSoup  # type: ignore
    except ImportError:
        raise ValueError("ebooklib and beautifulsoup4 are required for EPUB")

    import tempfile
    import os

    # ebooklib needs a file path
    with tempfile.NamedTemporaryFile(suffix=".epub", delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        book = epub.read_epub(tmp_path, options={"ignore_ncx": True})
        parts: list[str] = []

        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n")
            clean = "\n".join(line.strip() for line in text.splitlines() if line.strip())
            if clean:
                parts.append(clean)

        if not parts:
            raise ValueError("EPUB contained no extractable text")
        return "\n\n".join(parts)
    finally:
        os.unlink(tmp_path)


# ── EML (Email) ────────────────────────────────────────────────────


def _extract_eml(data: bytes) -> str:
    """Extract text from an EML email file."""
    import email
    from email import policy

    msg = email.message_from_bytes(data, policy=policy.default)

    parts: list[str] = []

    # Headers
    for header in ("From", "To", "Cc", "Subject", "Date"):
        val = msg.get(header)
        if val:
            parts.append(f"{header}: {val}")

    parts.append("")  # blank line

    # Body
    body = msg.get_body(preferencelist=("plain", "html"))
    if body:
        content = body.get_content()
        if body.get_content_type() == "text/html":
            try:
                from bs4 import BeautifulSoup  # type: ignore
                content = BeautifulSoup(content, "html.parser").get_text(separator="\n")
            except ImportError:
                import re
                content = re.sub(r"<[^>]+>", " ", content)

        parts.append(content.strip())

    text = "\n".join(parts)
    if not text.strip():
        raise ValueError("EML contained no extractable text")
    return text


# ── MSG (Outlook) ──────────────────────────────────────────────────


def _extract_msg(data: bytes) -> str:
    """Extract text from an Outlook MSG file."""
    try:
        import extract_msg  # type: ignore
    except ImportError:
        raise ValueError("extract-msg is not installed")

    import tempfile
    import os

    with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        msg = extract_msg.Message(tmp_path)
        parts: list[str] = []

        if msg.sender:
            parts.append(f"From: {msg.sender}")
        if msg.to:
            parts.append(f"To: {msg.to}")
        if msg.subject:
            parts.append(f"Subject: {msg.subject}")
        if msg.date:
            parts.append(f"Date: {msg.date}")

        parts.append("")

        if msg.body:
            parts.append(msg.body)

        msg.close()

        text = "\n".join(parts)
        if not text.strip():
            raise ValueError("MSG contained no extractable text")
        return text
    finally:
        os.unlink(tmp_path)


# ── ODT (OpenDocument) ─────────────────────────────────────────────


def _extract_odt(data: bytes) -> str:
    """Extract text from an ODT (OpenDocument Text) file."""
    try:
        from odf import text as odf_text  # type: ignore
        from odf import load as odf_load  # type: ignore
        from odf.text import P  # type: ignore
    except ImportError:
        raise ValueError("odfpy is not installed")

    doc = odf_load(io.BytesIO(data))
    paragraphs: list[str] = []

    for p in doc.getElementsByType(P):
        t = ""
        for node in p.childNodes:
            if hasattr(node, "data"):
                t += node.data
            elif hasattr(node, "__str__"):
                t += str(node)
        if t.strip():
            paragraphs.append(t.strip())

    if not paragraphs:
        raise ValueError("ODT contained no extractable text")
    return "\n\n".join(paragraphs)


# ── YAML ───────────────────────────────────────────────────────────


def _extract_yaml(data: bytes) -> str:
    """Extract readable text from YAML."""
    try:
        import yaml  # type: ignore
    except ImportError:
        # Fallback: just return as plain text
        return _extract_plaintext(data)

    text = _extract_plaintext(data)
    try:
        obj = yaml.safe_load(text)
        if isinstance(obj, (dict, list)):
            strings: list[str] = []
            _collect_strings_recursive(obj, strings)
            if strings:
                return "\n".join(strings)
        return text
    except Exception:
        return text


# ── TOML ───────────────────────────────────────────────────────────


def _extract_toml(data: bytes) -> str:
    """Extract readable text from TOML."""
    try:
        import tomllib  # Python 3.11+
    except ImportError:
        return _extract_plaintext(data)

    try:
        obj = tomllib.loads(_extract_plaintext(data))
        strings: list[str] = []
        _collect_strings_recursive(obj, strings)
        if strings:
            return "\n".join(strings)
        return _extract_plaintext(data)
    except Exception:
        return _extract_plaintext(data)


def _collect_strings_recursive(obj, out: list[str]) -> None:
    """Recursively collect string values from nested dicts/lists."""
    if isinstance(obj, str):
        if len(obj.strip()) > 3:
            out.append(obj.strip())
    elif isinstance(obj, dict):
        for k, v in obj.items():
            if isinstance(v, str) and len(v.strip()) > 3:
                out.append(f"{k}: {v.strip()}")
            else:
                _collect_strings_recursive(v, out)
    elif isinstance(obj, (list, tuple)):
        for item in obj:
            _collect_strings_recursive(item, out)


# ── Images (OCR via Tesseract) ─────────────────────────────────────


def _extract_image_ocr(data: bytes) -> str:
    """Extract text from images using Tesseract OCR."""
    try:
        from PIL import Image  # type: ignore
        import pytesseract  # type: ignore
    except ImportError:
        raise ValueError("Pillow and pytesseract are required for OCR. Install: pip install Pillow pytesseract")

    img = Image.open(io.BytesIO(data))
    text = pytesseract.image_to_string(img)

    if not text.strip():
        raise ValueError("Image contained no extractable text (OCR found nothing)")
    return text.strip()


# ── Audio (Whisper via Ollama or local) ────────────────────────────


def _extract_audio(data: bytes, mimetype: str) -> str:
    """Transcribe audio using OpenAI Whisper (runs locally on GPU)."""
    import tempfile
    import os

    ext_map = {
        "audio/mpeg": ".mp3",
        "audio/wav": ".wav",
        "audio/ogg": ".ogg",
        "audio/flac": ".flac",
        "audio/mp4": ".m4a",
        "audio/webm": ".webm",
    }
    ext = ext_map.get(mimetype, ".wav")

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        from faster_whisper import WhisperModel  # type: ignore

        # Use "base" model by default — good balance of speed and accuracy
        # Can be overridden via WHISPER_MODEL env var
        model_name = os.environ.get("WHISPER_MODEL", "base")
        device = "cuda" if os.environ.get("NVIDIA_VISIBLE_DEVICES") else "cpu"
        compute_type = "float16" if device == "cuda" else "int8"

        logger.info(f"Whisper: loading model '{model_name}' on {device} for audio transcription")
        model = WhisperModel(model_name, device=device, compute_type=compute_type)
        segments, info = model.transcribe(tmp_path, beam_size=5)

        text_parts: list[str] = []
        for segment in segments:
            text_parts.append(segment.text.strip())

        text = " ".join(text_parts)

        if not text.strip():
            raise ValueError("Audio contained no recognizable speech")

        logger.info(f"Whisper: transcribed {len(text)} chars, language: {info.language} ({info.language_probability:.0%})")
        return text.strip()
    except ImportError:
        raise ValueError("faster-whisper is not installed")
    finally:
        os.unlink(tmp_path)


def _extract_video(data: bytes, mimetype: str) -> str:
    """Extract audio track from video using ffmpeg, then transcribe with Whisper."""
    import tempfile
    import os
    import subprocess

    ext_map = {
        "video/mp4": ".mp4",
        "video/webm": ".webm",
        "video/quicktime": ".mov",
        "video/x-msvideo": ".avi",
        "video/x-matroska": ".mkv",
        "video/mpeg": ".mpeg",
    }
    ext = ext_map.get(mimetype, ".mp4")

    # Write video to temp file
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        video_path = tmp.name

    audio_path = video_path + ".wav"

    try:
        # Extract audio track with ffmpeg
        logger.info(f"ffmpeg: extracting audio from video ({mimetype})")
        result = subprocess.run(
            ["ffmpeg", "-i", video_path, "-vn", "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1", audio_path, "-y"],
            capture_output=True, text=True, timeout=300,
        )
        if result.returncode != 0:
            raise ValueError(f"ffmpeg failed to extract audio: {result.stderr[:200]}")

        if not os.path.exists(audio_path) or os.path.getsize(audio_path) < 100:
            raise ValueError("Video contained no extractable audio track")

        # Read the extracted audio and transcribe
        with open(audio_path, "rb") as f:
            audio_data = f.read()

        return _extract_audio(audio_data, "audio/wav")
    finally:
        for path in (video_path, audio_path):
            try:
                os.unlink(path)
            except OSError:
                pass
